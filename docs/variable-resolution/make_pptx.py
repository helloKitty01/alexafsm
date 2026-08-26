#!/usr/bin/env python3
"""生成 工具参数变量解析机制.pptx（浅色系，单页，与 .light.slides.html 内容一致）。"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 配色（浅色系，与 reminder 幻灯片一致）----
TEXT = RGBColor(0x1E, 0x2A, 0x3D)
MUTED = RGBColor(0x5A, 0x6B, 0x84)
ACCENT = RGBColor(0x2F, 0x6F, 0xED)
ACCENT2 = RGBColor(0x0E, 0xA5, 0x78)
WARN = RGBColor(0xD9, 0x7A, 0x06)
PURPLE = RGBColor(0x7C, 0x5C, 0xD6)
CARD_BORDER = RGBColor(0xDF, 0xE6, 0xF0)
CODE_BG = RGBColor(0xF3, 0xF6, 0xFB)
CODE_TEXT = RGBColor(0x24, 0x40, 0x5F)
SLIDE_BG = RGBColor(0xF8, 0xFB, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEAD_BG = RGBColor(0xEC, 0xF2, 0xFC)

FONT = "Microsoft YaHei"
MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def _set_ea_font(run):
    rPr = run.font._rPr
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def add_card(slide, x, y, w, h, fill=WHITE, border=CARD_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.05
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(1)
    card.shadow.inherit = False
    return card


def text_box(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.12, space_after=3):
    """lines: list of paragraphs; each paragraph is list of (text, size, color, bold, mono)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for text, size, color, bold, mono in line:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = MONO if mono else FONT
            if not mono:
                _set_ea_font(r)
    return tb


def code_box(slide, x, y, w, h, lines, size=8.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.08
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = CARD_BORDER
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(7)
    tf.margin_top = tf.margin_bottom = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.1
        for text, color, bold in line:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = MONO
            _set_ea_font(r)
    return box


s = prs.slides.add_slide(blank)
s.background.fill.solid()
s.background.fill.fore_color.rgb = SLIDE_BG
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Pt(4))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background()

M = Inches(0.42)

# ---- 标题 + 语法条 ----
text_box(s, M, Inches(0.18), Inches(4.5), Inches(0.6),
         [[("工具参数变量解析机制方案", 20, TEXT, True, False)]],
         anchor=MSO_ANCHOR.MIDDLE)
add_card(s, Inches(5.0), Inches(0.16), Inches(7.92), Inches(0.78), fill=HEAD_BG)
text_box(s, Inches(5.2), Inches(0.22), Inches(7.55), Inches(0.68), [
    [("${ app . 抖音 . bundleName }", 12, ACCENT, True, True)],
    [("namespace . entity . attribute 三段式，命名空间可扩展；模型只表达\u201c哪个应用\u201d，真实标识符由框架在工具调用前查表求值。", 9, MUTED, False, False)],
], anchor=MSO_ANCHOR.MIDDLE, space_after=2)

TOP = Inches(1.10)
CARD_H = Inches(3.45)
LW = Inches(4.55)          # 左卡宽
RX = M + LW + Inches(0.22) # 右卡 x
RW = SLIDE_W - RX - M      # 右卡宽

# ---- 左卡：问题与目标 ----
add_card(s, M, TOP, LW, CARD_H)
text_box(s, M + Inches(0.2), TOP + Inches(0.12), LW - Inches(0.4), CARD_H - Inches(0.24), [
    [("● 问题", 11.5, WARN, True, False)],
    [("• 两次 loop：", 9, TEXT, True, False),
     ("模型须先调 getAllInstalledApps 查 bundleName，再调 forbidPermission(bundleName, 权限)", 9, MUTED, False, False)],
    [("• 幻觉风险：", 9, TEXT, True, False),
     ("凭记忆填写 bundleName 可能编造（如 com.douyin.app）", 9, MUTED, False, False)],
    [("• 硬约束：", 9, TEXT, True, False),
     ("工具不可修改，带 bundleName 的工具很多，无法逐个改造", 9, MUTED, False, False)],
    [("", 6, MUTED, False, False)],
    [("● 目标", 11.5, ACCENT2, True, False)],
    [("• 正常路径 ", 9, MUTED, False, False),
     ("1 次 loop", 9, TEXT, True, False),
     (" 完成调用", 9, MUTED, False, False)],
    [("• 工具零改动", 9, TEXT, True, False),
     ("，同类工具自动全覆盖", 9, MUTED, False, False)],
    [("• 真实标识符永远来自查表", 9, TEXT, True, False),
     ("，结构性消除幻觉", 9, MUTED, False, False)],
], space_after=5)

# ---- 右卡：关键设计 ----
add_card(s, RX, TOP, RW, CARD_H)
PX = RX + Inches(0.2)
PW = RW - Inches(0.4)
text_box(s, PX, TOP + Inches(0.12), PW, Inches(0.3),
         [[("● 关键设计", 11.5, ACCENT, True, False)]])

# 两张表
text_box(s, PX, TOP + Inches(0.48), Inches(1.35), Inches(0.6), [
    [("两张表", 9.5, TEXT, True, False)],
    [("数据基础", 8, ACCENT, False, False)],
], space_after=1)
half = (PW - Inches(1.45) - Inches(0.15)) / 2
code_box(s, PX + Inches(1.45), TOP + Inches(0.44), half, Inches(0.68), [
    [("同义词表 · 别名归一", TEXT, True)],
    [("douyin / Douyin → 抖音", CODE_TEXT, False)],
], size=8.5)
code_box(s, PX + Inches(1.45) + half + Inches(0.15), TOP + Inches(0.44), half, Inches(0.68), [
    [("映射表 · 规范名 → bundleName", TEXT, True)],
    [("抖音 → com.ss.android.ugc.aweme", CODE_TEXT, False)],
], size=8.5)

# Pre-call 拦截器
text_box(s, PX, TOP + Inches(1.36), Inches(1.35), Inches(0.8), [
    [("Pre-call", 9.5, TEXT, True, False)],
    [("拦截器", 9.5, TEXT, True, False)],
    [("框架层·工具无感知", 7.5, ACCENT, False, False)],
], space_after=1)
text_box(s, PX + Inches(1.45), TOP + Inches(1.34), PW - Inches(1.45), Inches(0.55), [
    [("工具调用前扫描参数中的 ${...}，查表替换为真实值；解析失败则短路，不把变量字面量透传给工具。", 9, MUTED, False, False)],
])
code_box(s, PX + Inches(1.45), TOP + Inches(1.92), PW - Inches(1.45), Inches(0.42), [
    [("扫描 ${...}  →  同义词归一  →  查映射表  →  替换真实值", CODE_TEXT, True)],
], size=9)

# Skill 设计
text_box(s, PX, TOP + Inches(2.52), Inches(1.35), Inches(0.6), [
    [("Skill 设计", 9.5, TEXT, True, False)],
    [("提示词约定", 8, ACCENT, False, False)],
], space_after=1)
text_box(s, PX + Inches(1.45), TOP + Inches(2.50), PW - Inches(1.45), Inches(0.5), [
    [("bundleName 参数直接填 ${app.应用名.bundleName}，无需先查已安装列表，禁止凭记忆编造标识符。", 9, MUTED, False, False)],
])
code_box(s, PX + Inches(1.45), TOP + Inches(2.94), PW - Inches(1.45), Inches(0.4), [
    [("forbidPermission(bundleName=", CODE_TEXT, False),
     ('"${app.抖音.bundleName}"', ACCENT2, True),
     (", permission=", CODE_TEXT, False),
     ('"麦克风"', ACCENT2, True),
     (")", CODE_TEXT, False)],
], size=8.5)

# ---- 底部横条：关键流程 ----
BY = TOP + CARD_H + Inches(0.2)
BH = SLIDE_H - BY - Inches(0.30)
add_card(s, M, BY, SLIDE_W - M * 2, BH)
BW = SLIDE_W - M * 2
half_flow = (BW - Inches(0.6)) / 2

text_box(s, M + Inches(0.2), BY + Inches(0.10), half_flow, Inches(0.28),
         [[("正常路径 · 1 次 loop", 9.5, ACCENT2, True, False)]])
code_box(s, M + Inches(0.2), BY + Inches(0.40), half_flow - Inches(0.1), Inches(0.55), [
    [("模型输出 ${app.抖音.bundleName} → 拦截器查表（同义词+映射）", CODE_TEXT, False)],
    [("→ 替换真实值 com.ss.android... → 执行工具", CODE_TEXT, False)],
], size=8.5)

FX2 = M + Inches(0.2) + half_flow + Inches(0.2)
text_box(s, FX2, BY + Inches(0.10), half_flow, Inches(0.28),
         [[("兜底路径 · 解析失败时（可选）", 9.5, WARN, True, False)]])
code_box(s, FX2, BY + Inches(0.40), half_flow - Inches(0.1), Inches(0.55), [
    [("解析失败 → 查询全量已安装应用 getAllInstalledApps", CODE_TEXT, False)],
    [("→ 列表 + 错误作为 tool result 给模型 → 模型纠正后重试", CODE_TEXT, False)],
], size=8.5)

text_box(s, M + Inches(0.2), BY + Inches(1.08), BW - Inches(0.4), Inches(0.75), [
    [("进一步优化  ", 9, PURPLE, True, False),
     ("① 标准名列表注入：", 9, TEXT, True, False),
     ("会话开始时把已安装应用的标准名称列表给到模型，从封闭词表选名，兜底触发概率大幅降低。  ", 9, MUTED, False, False),
     ("② 工具 schema 改写：", 9, TEXT, True, False),
     ("框架注册层把 bundleName 参数重定义为 appName，拦截器转换后调用原工具——模型直接传\u201c抖音\u201d；变量语法保留为通用机制。", 9, MUTED, False, False)],
    [("◎ 方案不变量：", 9, ACCENT2, True, False),
     ("模型永远只表达\u201c哪个应用\u201d，真实标识符永远来自查表 —— 幻觉被结构性消除，绝大多数情况 loop 压到 1 次。", 9, TEXT, False, False)],
], space_after=4)

out = "/workspace/docs/variable-resolution/工具参数变量解析机制.pptx"
prs.save(out)
print("saved:", out)
