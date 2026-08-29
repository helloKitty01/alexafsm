#!/usr/bin/env python3
"""生成浅色系汇报 PPT（.pptx，16:9，两页：现状与问题 / 目标方案）。

用法：python3 make_pptx.py
输出：上下文工程与KV-Cache优化.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 浅色主题色板 ----
TEXT   = RGBColor(0x1F, 0x29, 0x37)   # 主文字
MUTED  = RGBColor(0x64, 0x74, 0x8B)   # 次要文字
BLUE   = RGBColor(0x25, 0x63, 0xEB)   # 强调蓝
GREEN  = RGBColor(0x05, 0x96, 0x69)   # 正向绿
RED    = RGBColor(0xDC, 0x26, 0x26)   # 问题红
AMBER  = RGBColor(0xD9, 0x77, 0x06)   # 警示橙
CARD   = RGBColor(0xF8, 0xFA, 0xFC)   # 卡片底
BORDER = RGBColor(0xE2, 0xE8, 0xF0)   # 卡片边
RED_BG    = RGBColor(0xFE, 0xE2, 0xE2)
RED_BD    = RGBColor(0xFC, 0xA5, 0xA5)
RED_TX    = RGBColor(0xB9, 0x1C, 0x1C)
GREEN_BG  = RGBColor(0xD1, 0xFA, 0xE5)
GREEN_BD  = RGBColor(0x6E, 0xE7, 0xB7)
GREEN_TX  = RGBColor(0x04, 0x78, 0x57)
BLUE_BG   = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_BD   = RGBColor(0x93, 0xC5, 0xFD)
AMBER_BG  = RGBColor(0xFE, 0xF3, 0xC7)
AMBER_BD  = RGBColor(0xFC, 0xD3, 0x4D)
CODE_BG   = RGBColor(0xF1, 0xF5, 0xF9)

FONT = "Microsoft YaHei"


def set_font(run, size, color=TEXT, bold=False, mono=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    name = "Consolas" if mono else FONT
    f.name = name
    # 同时设置中文字体（a:ea）
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


def box(slide, x, y, w, h, fill=CARD, border=BORDER, rounded=True, dash=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            shape.adjustments[0] = 0.12
        except Exception:
            pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1)
    if dash:
        ln = shape.line._get_or_add_ln()
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(d)
    shape.shadow.inherit = False
    return shape


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.0, wrap=True):
    """runs: list[list[(text, kwargs)]]，外层每项为一段。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for t, kw in para:
            r = p.add_run()
            r.text = t
            set_font(r, **kw)
    return tb


def vband(slide, x, y, w, h, label, bg, bd, tx):
    shape = box(slide, x, y, w, h, fill=bg, border=bd)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    set_font(r, 8, color=tx, bold=True)
    # 竖排文字
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("vert", "eaVert")
    return shape


def kpi_card(slide, x, y, w, h, big, big_color, label):
    box(slide, x, y, w, h)
    text(slide, x + 0.1, y + 0.07, w - 0.2, 0.32,
         [[(big, dict(size=16, color=big_color, bold=True))]],
         align=PP_ALIGN.CENTER)
    text(slide, x + 0.1, y + 0.4, w - 0.2, h - 0.45,
         [[(label, dict(size=8.5, color=MUTED))]],
         align=PP_ALIGN.CENTER, line_spacing=1.05)


def layer_row(slide, x, y, w, h, name, name_color, desc_runs, fill, bd, dash=None):
    box(slide, x, y, w, h, fill=fill, border=bd, dash=dash)
    text(slide, x + 0.12, y, 1.15, h, [[(name, dict(size=9.5, color=name_color, bold=True))]],
         anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    text(slide, x + 1.32, y, w - 1.45, h, [desc_runs],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)


def ent_box(slide, x, y, w, h, title, sub):
    box(slide, x, y, w, h, fill=RGBColor(0xFF, 0xFF, 0xFF), border=MUTED, dash="dash")
    text(slide, x + 0.05, y, w - 0.1, h,
         [[(title, dict(size=9, color=TEXT, bold=True))],
          [(sub, dict(size=7.5, color=MUTED))]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)


def banner(slide, x, y, w, h, runs, bg, bd):
    box(slide, x, y, w, h, fill=bg, border=bd)
    text(slide, x + 0.25, y, w - 0.5, h, [runs], anchor=MSO_ANCHOR.MIDDLE)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ============================================================
# 第 1 页：现状与问题
# ============================================================
s = prs.slides.add_slide(blank)

text(s, 0.35, 0.18, 12.6, 0.45,
     [[("现状与问题：", dict(size=20, bold=True)),
       ("命中率为何只有 66%", dict(size=20, bold=True, color=BLUE)),
       (" —— 失效范围含 Tools 与全部历史", dict(size=20, bold=True))]])

# KPI 条
kw = (12.63 - 0.4) / 3
kpi_card(s, 0.35, 0.72, kw, 0.78, "66%", AMBER, "现网 KV cache 命中率")
kpi_card(s, 0.35 + kw + 0.2, 0.72, kw, 0.78, "37K / 25K", BLUE, "平均每请求输入 / 命中（miss 12K）")
kpi_card(s, 0.35 + 2 * (kw + 0.2), 0.72, kw, 0.78, "~60K", BLUE, "会话上下文规模")

# 左：结构图卡
box(s, 0.35, 1.62, 7.15, 4.92)
text(s, 0.55, 1.72, 6.8, 0.3,
     [[("三个独立实体的物理拼接（自上而下 = 前缀从前到后）", dict(size=11, bold=True))]])

rx, rw = 1.72, 4.98          # 层级行的 x / 宽
ey, step, rh = 2.12, 0.60, 0.52
# 实体列
ent_box(s, 0.52, ey, 1.1, 3 * step - 0.08, "① System Prompt", "单字段：产品/会话/query级拼于正文内")
ent_box(s, 0.52, ey + 3 * step, 1.1, rh, "② Tools", "独立字段")
ent_box(s, 0.52, ey + 4 * step, 1.1, 3 * step - 0.08, "③ UAT", "messages 消息流（loop级）")
# 层级行
layer_row(s, rx, ey, rw, rh, "产品级", GREEN,
          [("agent 关键原则，所有用户相同 ✓", dict(size=8.5, color=MUTED))], CARD, GREEN_BD)
layer_row(s, rx, ey + step, rw, rh, "会话级", GREEN,
          [("用户信息、长期记忆，多轮更新一次 ✓", dict(size=8.5, color=MUTED))], CARD, GREEN_BD)
layer_row(s, rx, ey + 2 * step, rw, rh, "query级", RED,
          [("available skills（基于 query 搜索）+ 相关知识/用户信息 + 环境上下文，", dict(size=8.5, color=MUTED)),
           ("每 query 原位重写 ✗", dict(size=8.5, color=RED, bold=True))], RED_BG, RED_BD)
layer_row(s, rx, ey + 3 * step, rw, rh, "Tools块", AMBER,
          [("内容稳定，但模板拼接于 SP 之后 → ", dict(size=8.5, color=MUTED)),
           ("被 query级重写连带 ✗", dict(size=8.5, color=RED, bold=True))], AMBER_BG, AMBER_BD)
layer_row(s, rx, ey + 4 * step, rw, rh, "loop 1（最旧）", BLUE,
          [("U→[A⇄T ×n]→A　", dict(size=8.5, color=MUTED)),
           ("⚡10 轮滑窗从此裁剪，每轮再断一次", dict(size=8.5, color=RED, bold=True))], BLUE_BG, BLUE_BD)
layer_row(s, rx, ey + 5 * step, rw, rh, "loop 2", BLUE,
          [("U→[A⇄T ×n]→A，loop 内迭代可命中", dict(size=8.5, color=MUTED))], BLUE_BG, BLUE_BD)
layer_row(s, rx, ey + 6 * step, rw, rh, "loop 3 …", BLUE,
          [("追加中", dict(size=8.5, color=MUTED))], CARD, BLUE_BD, dash="dash")
# 红色失效竖带（从 query级行起贯穿到底）
vband(s, 6.82, ey + 2 * step, 0.5, 5 * step - 0.08,
      "⚡ 每 query 从此失效：query级 + Tools + UAT 全部历史", RED_BG, RED_BD, RED_TX)
# loop 精简说明
text(s, 0.55, ey + 7 * step + 0.02, 6.8, 0.3,
     [[("loop 内部：U query → A call ⇄ T result ×n → A 答复　", dict(size=8.5, color=MUTED)),
       ("新 query（loop 切换）= miss 点", dict(size=8.5, color=RED, bold=True)),
       ("｜", dict(size=8.5, color=MUTED)),
       ("loop 内迭代 = 追加命中", dict(size=8.5, color=GREEN, bold=True))]])

# 右：四个问题卡
px, pw = 7.65, 5.33
ph, pgap = 1.14, 0.12
problems = [
    ("① query 切换点全量 miss",
     [("query级内容重写 → 失效 = ", dict(size=8.8, color=MUTED)),
      ("query级 + Tools + UAT 全部历史", dict(size=8.8, color=TEXT, bold=True)),
      ("，随轮数线性增长；Tools 内容稳定却被连带重算。loop 内可命中，损失集中在每个新 query 首轮。", dict(size=8.8, color=MUTED))]),
    ("② 10 轮滑窗裁剪 = 长对话每轮全量 miss",
     [("聊过 10 轮后", dict(size=8.8, color=MUTED)),
      ("每新一轮裁最旧一轮", dict(size=8.8, color=TEXT, bold=True)),
      ("，前缀头部每轮变化；T 体积大（skill 正文）被单独高频裁剪。", dict(size=8.8, color=MUTED)),
      ("长对话阶段比①更严重。", dict(size=8.8, color=TEXT, bold=True))]),
    ("③ 相关性绑定丢失 + 历史证据被覆盖",
     [("query级槽位与 query ", dict(size=8.8, color=MUTED)),
      ("无显式对应", dict(size=8.8, color=TEXT, bold=True)),
      ("；每轮覆盖式更新，第 N 轮依赖第 N-2 轮检索结果时证据已丢，", dict(size=8.8, color=MUTED)),
      ("去重无从谈起。", dict(size=8.8, color=TEXT, bold=True))]),
    ("④ 驻留竞争（eviction）",
     [("自建集群缓存容量有限，被", dict(size=8.8, color=MUTED)),
      ("全网用户流量逐出", dict(size=8.8, color=TEXT, bold=True)),
      ("——前缀不变 ≠ 命中。排布性 miss 与驻留性 miss 需分开度量、分开治理。", dict(size=8.8, color=MUTED))]),
]
for i, (title, body) in enumerate(problems):
    y = 1.62 + i * (ph + pgap)
    box(s, px, y, pw, ph)
    text(s, px + 0.18, y + 0.09, pw - 0.36, 0.28,
         [[(title, dict(size=10, color=RED, bold=True))]])
    text(s, px + 0.18, y + 0.38, pw - 0.36, ph - 0.45, [body], line_spacing=1.05)

# 底部结论 banner
banner(s, 0.35, 6.72, 12.63, 0.5,
       [("综上，缓存失效有两个来源：", dict(size=10, color=TEXT)),
        ("问题① query级内容每次重写", dict(size=10, color=RED, bold=True)),
        (" ＋ ", dict(size=10, color=TEXT)),
        ("问题② 滑窗裁剪每轮改历史", dict(size=10, color=RED, bold=True)),
        ("；另有问题④ 驻留竞争需集群侧配合 —— 下页逐项治理", dict(size=10, color=TEXT))],
       RED_BG, RED_BD)

# ============================================================
# 第 2 页：目标方案
# ============================================================
s = prs.slides.add_slide(blank)

text(s, 0.35, 0.15, 12.6, 0.42,
     [[("目标方案：上下文只增不改，", dict(size=20, bold=True)),
       ("挑战 90%+ 命中率", dict(size=20, bold=True, color=BLUE))]])
text(s, 0.35, 0.6, 12.6, 0.28,
     [[("与现状图同构对照：", dict(size=9.5, color=MUTED)),
       ("治问题①③", dict(size=9.5, color=RED, bold=True)),
       ("（query级下沉，两条通道）· ", dict(size=9.5, color=MUTED)),
       ("治问题②", dict(size=9.5, color=RED, bold=True)),
       ("（裁剪改造）· ", dict(size=9.5, color=MUTED)),
       ("答问题④", dict(size=9.5, color=RED, bold=True)),
       ("（驻留诉求）", dict(size=9.5, color=MUTED))]])

# 左：同构结构图
box(s, 0.35, 0.98, 7.15, 4.28)
text(s, 0.55, 1.07, 6.8, 0.28,
     [[("同一物理序列，改动点逐行标注", dict(size=11, bold=True))]])
rx, rw = 1.72, 4.98
ey, step, rh = 1.42, 0.58, 0.5
ent_box(s, 0.52, ey, 1.1, 3 * step - 0.08, "① System Prompt", "仅产品级+会话级，会话内字节不变")
ent_box(s, 0.52, ey + 3 * step, 1.1, rh, "② Tools", "独立字段")
ent_box(s, 0.52, ey + 4 * step, 1.1, 2 * step - 0.08, "③ UAT", "messages 消息流（loop级）")

rows2 = [
    ("产品级", GREEN, [("agent 关键原则", dict(size=8.5, color=MUTED)),
                       ("　【不变】", dict(size=8, color=GREEN, bold=True))], CARD, GREEN_BD, None),
    ("会话级", GREEN, [("用户信息、长期记忆，会话冻结、更新攒批", dict(size=8.5, color=MUTED)),
                       ("　【策略化】", dict(size=8, color=GREEN, bold=True))], CARD, GREEN_BD, None),
    ("query级", RED, [("槽位已取消", dict(size=8.5, color=TEXT, bold=True)),
                      (" → 内容经两条通道下沉至 UAT（见右侧）", dict(size=8.5, color=MUTED)),
                      ("　【取消 · 治①③】", dict(size=8, color=RED, bold=True))], RGBColor(0xFF, 0xFF, 0xFF), RED_BD, "dash"),
    ("Tools块", AMBER, [("前方不再有改写位置，", dict(size=8.5, color=MUTED)),
                        ("稳定命中", dict(size=8.5, color=TEXT, bold=True)),
                        ("　【被解救 ✓】", dict(size=8, color=GREEN, bold=True))], AMBER_BG, AMBER_BD, None),
    ("已裁剪区", GREEN, [("按长度水位、一次裁足，裁剪由", dict(size=8.5, color=MUTED)),
                         ("每轮降为几十轮一次", dict(size=8.5, color=TEXT, bold=True)),
                         ("　【改造 · 治②】", dict(size=8, color=BLUE, bold=True))], GREEN_BG, GREEN_BD, "dash"),
    ("loop原文＋追加", BLUE, [("U（turn_context）→[A⇄T（search/load）]→A，仅为增量付 prefill", dict(size=8.5, color=MUTED)),
                              ("　【注入增强】", dict(size=8, color=BLUE, bold=True))], BLUE_BG, BLUE_BD, None),
]
for i, (name, nc, desc, fill, bd, dash) in enumerate(rows2):
    layer_row(s, rx, ey + i * step, rw, rh, name, nc, desc, fill, bd, dash=dash)
vband(s, 6.82, ey, 0.5, 6 * step - 0.08,
      "✓ 每轮前缀 = 上一轮严格超集，排布性 miss 归零", GREEN_BG, GREEN_BD, GREEN_TX)
text(s, 0.55, ey + 6 * step + 0.02, 6.8, 0.44,
     [[("对照上页现状图逐行阅读：query级槽位取消后，System Prompt 与 Tools 在会话内完全静态；"
        "UAT 只增不改，动态内容全部经右侧两条通道进入消息流。", dict(size=8.5, color=MUTED))]],
     line_spacing=1.05)

# 右：通道一 / 通道二
px, pw = 7.65, 5.33
box(s, px, 0.98, pw, 2.62)
text(s, px + 0.18, 1.07, pw - 0.36, 0.28,
     [[("通道一 · 被动注入：", dict(size=10.5, bold=True)),
       ("turn_context 随 U 进入", dict(size=10.5, bold=True, color=BLUE)),
       ("　治①③", dict(size=8.5, color=RED, bold=True))]])
box(s, px + 0.18, 1.4, pw - 0.36, 1.22, fill=CODE_BG, border=BORDER, rounded=False)
code_lines = [
    [("<turn_context>", dict(size=8, color=MUTED, mono=True))],
    [("  <env>", dict(size=8, color=BLUE, mono=True)),
     ("环境增量（首轮全量快照，此后只注增量）", dict(size=8, color=TEXT, mono=True)),
     ("</env>", dict(size=8, color=BLUE, mono=True))],
    [("  <retrieved>", dict(size=8, color=GREEN, mono=True)),
     ("相关知识/用户信息（ID 去重，只注新条目）", dict(size=8, color=TEXT, mono=True)),
     ("</retrieved>", dict(size=8, color=GREEN, mono=True))],
    [("  <skills>", dict(size=8, color=AMBER, mono=True)),
     ("available skills（方案 A 的注入位置）", dict(size=8, color=TEXT, mono=True)),
     ("</skills>", dict(size=8, color=AMBER, mono=True))],
    [("</turn_context>", dict(size=8, color=MUTED, mono=True))],
    [("<query>", dict(size=8, color=BLUE, mono=True)),
     ("用户原文，永远置于消息最后", dict(size=8, color=TEXT, mono=True)),
     ("</query>", dict(size=8, color=BLUE, mono=True))],
]
text(s, px + 0.3, 1.46, pw - 0.6, 1.12, code_lines, line_spacing=1.1)
text(s, px + 0.18, 2.68, pw - 0.36, 0.86,
     [[("· 拼接结果逐字节稳定：标签顺序/条目排序固定、空块不输出、发出后不再修改（否则缓存悄悄失效）",
        dict(size=8.4, color=MUTED))],
      [("· 检索内容与 query ", dict(size=8.4, color=MUTED)),
       ("同消息相邻、永久留在历史", dict(size=8.4, color=TEXT, bold=True)),
       ("：对应关系模型可见，可去重、可审计（治问题③）", dict(size=8.4, color=MUTED))]],
     line_spacing=1.12)

box(s, px, 3.68, pw, 1.58)
text(s, px + 0.18, 3.77, pw - 0.36, 0.28,
     [[("通道二 · 主动加载：", dict(size=10.5, bold=True)),
       ("skills 以 T 进入 loop", dict(size=10.5, bold=True, color=GREEN)),
       ("　治①", dict(size=8.5, color=RED, bold=True))]])
text(s, px + 0.18, 4.08, pw - 0.36, 0.26,
     [[("search_skill(query) → T: available skills → load_skill(name) → T: 正文（load 沿用现状）",
        dict(size=8.3, color=TEXT, mono=True))]])
text(s, px + 0.18, 4.36, pw - 0.36, 0.86,
     [[("· 方案 A（先行）：保留系统侧 skill search，结果随 U 注入，只改注入位置、检索链路不动",
        dict(size=8.8, color=MUTED))],
      [("· 方案 B（演进）：search 工具化，模型自主决定检索时机与关键词，可按阶段多次检索",
        dict(size=8.8, color=MUTED))]], line_spacing=1.1)

# 第二行：裁剪 / 诉求
half = (12.63 - 0.2) / 2
box(s, 0.35, 5.38, half, 0.86)
text(s, 0.53, 5.46, half - 0.36, 0.24,
     [[("裁剪机制改造", dict(size=10, bold=True)), ("　治②", dict(size=8.5, color=RED, bold=True))]])
text(s, 0.53, 5.72, half - 0.36, 0.48,
     [[("触发：按轮数（10 轮）→ 按 token 长度水位；比例：一次裁足到低水位（80% 触发、裁至 40%）；"
        "T 与 UAT 攒批同刻处理。效果：裁剪从每轮降为几十轮一次", dict(size=8.6, color=MUTED))]],
     line_spacing=1.05)
box(s, 0.35 + half + 0.2, 5.38, half, 0.86)
text(s, 0.53 + half + 0.2, 5.46, half - 0.36, 0.24,
     [[("对推理集群的诉求", dict(size=10, bold=True)), ("　答④", dict(size=8.5, color=RED, bold=True))]])
text(s, 0.53 + half + 0.2, 5.72, half - 0.36, 0.48,
     [[("诉求一：同一多轮对话及 loop 内，KV cache 会话存续期内不被逐出；"
        "诉求二：命中率按排布性 / 驻留性拆分度量", dict(size=8.6, color=MUTED))]],
     line_spacing=1.05)

# 收益 KPI
kw = (12.63 - 0.4) / 3
kpi_card(s, 0.35, 6.32, kw, 0.72, "66% → ~92%", GREEN, "命中率：现网实测 → 目标预估（miss 12K → ~3K / 请求）")
kpi_card(s, 0.35 + kw + 0.2, 6.32, kw, 0.72, "-56%", GREEN, "每请求等效输入成本：14.5K → 6.4K 等效 token")
kpi_card(s, 0.35 + 2 * (kw + 0.2), 6.32, kw, 0.72, "-75%", GREEN, "每请求 prefill 计算量：12K → 3K，TTFT 与吞吐同步改善")

# 落地节奏
banner(s, 0.35, 7.1, 12.63, 0.34,
       [("落地节奏：① 命中率拆分度量 → ② 裁剪改造 → ③ skills 迁移 → ④ 结构化注入 → ⑤ 驻留治理",
         dict(size=9, color=TEXT, bold=True))],
       BLUE_BG, BLUE_BD)

prs.save("上下文工程与KV-Cache优化.pptx")
print("saved: 上下文工程与KV-Cache优化.pptx")
