#!/usr/bin/env python3
"""生成 reminder机制.pptx（浅色系，2 页，与 reminder机制.slides.html 内容一致）。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 配色（浅色系）----
TEXT = RGBColor(0x1E, 0x2A, 0x3D)
MUTED = RGBColor(0x5A, 0x6B, 0x84)
ACCENT = RGBColor(0x2F, 0x6F, 0xED)
ACCENT2 = RGBColor(0x0E, 0xA5, 0x78)
WARN = RGBColor(0xD9, 0x7A, 0x06)
CARD_BORDER = RGBColor(0xDF, 0xE6, 0xF0)
CODE_BG = RGBColor(0xF3, 0xF6, 0xFB)
CODE_TEXT = RGBColor(0x24, 0x40, 0x5F)
SLIDE_BG = RGBColor(0xF8, 0xFB, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEAD_BG = RGBColor(0xEC, 0xF2, 0xFC)

FONT = "Microsoft YaHei"
MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = SLIDE_BG


def add_topbar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_card(slide, x, y, w, h, fill=WHITE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.045
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(1)
    card.shadow.inherit = False
    return card


def _set_ea_font(run):
    rPr = run.font._rPr
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def text_box(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.12, space_after=3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for text, size, color, bold, mono in line:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = MONO if mono else FONT
            if not mono:
                _set_ea_font(r)
    return tb


def add_table(slide, x, y, w, h, headers, rows, col_widths, body_size=9.5):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h)
    tbl = shape.table
    # 关闭默认样式条纹
    tbl.first_row = False
    tbl.horz_banding = False
    for j, cw in enumerate(col_widths):
        tbl.columns[j].width = Emu(int(w * cw))
    for j, htxt in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEAD_BG
        cell.margin_left = cell.margin_right = Pt(5)
        cell.margin_top = cell.margin_bottom = Pt(2)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = htxt
        r.font.size = Pt(body_size)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        r.font.name = FONT
        _set_ea_font(r)
    for i, row in enumerate(rows):
        for j, cell_txt in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            cell.margin_left = cell.margin_right = Pt(5)
            cell.margin_top = cell.margin_bottom = Pt(2)
            p = cell.text_frame.paragraphs[0]
            p.line_spacing = 1.05
            r = p.add_run()
            r.text = cell_txt
            r.font.size = Pt(body_size)
            r.font.bold = (j == 0)
            r.font.color.rgb = TEXT if j == 0 else MUTED
            r.font.name = FONT
            _set_ea_font(r)
            cell.text_frame.word_wrap = True
    return shape


def add_code_block(slide, x, y, w, h, lines):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.06
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = CARD_BORDER
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(8)
    tf.margin_top = tf.margin_bottom = Pt(5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.12
        for text, color, bold in line:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(8.5)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = MONO
            _set_ea_font(r)
    return box


def add_pageno(slide, txt):
    text_box(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.38),
             Inches(0.8), Inches(0.3), [[(txt, 9, MUTED, False, False)]])


# ══════════════════ 第 1 页：业界洞察 ══════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1)
add_topbar(s1)

M = Inches(0.42)          # 页边距
CW = Inches(6.1)          # 半栏宽

# 标题
text_box(s1, M, Inches(0.18), Inches(4.4), Inches(0.75), [
    [("Reminder 机制 · 业界洞察", 21, TEXT, True, False)],
    [("内部代号：notion 方案", 10, MUTED, False, False)],
])
# 标题右侧：定义条
head = add_card(s1, Inches(4.95), Inches(0.16), Inches(7.97), Inches(0.78), fill=HEAD_BG)
text_box(s1, Inches(5.15), Inches(0.22), Inches(7.6), Inches(0.68), [
    [("定义：", 9.5, TEXT, True, False),
     ("agent loop 运行期间，由系统向消息流注入的短小、结构化指令/状态块 —— agent 的二级指令通道。", 9.5, MUTED, False, False)],
    [("system prompt 是入职培训（primacy，静态），reminder 是定期备忘录（recency，事件触发）。", 9.5, ACCENT2, True, False)],
], anchor=MSO_ANCHOR.MIDDLE, space_after=1)

TOP = Inches(1.08)
CARD_H = Inches(3.28)

# 左卡：业界三条路线
add_card(s1, M, TOP, CW, CARD_H)
text_box(s1, M + Inches(0.2), TOP + Inches(0.1), CW - Inches(0.4), Inches(0.3),
         [[("● 业界三条路线（按\u201c谁来提醒\u201d划分）", 11.5, TEXT, True, False)]])
add_table(s1, M + Inches(0.18), TOP + Inches(0.46), CW - Inches(0.36), Inches(1.7),
          ["路线", "代表", "机制"],
          [["系统注入", "Claude Code\n<system-reminder>",
            "40+ 种事件触发器生成提醒，附着到 user message / tool result，一条不进 system prompt"],
           ["自我复诵", "Manus recitation",
            "模型反复重写 todo.md 到上下文尾部，把全局计划推入高注意力区"],
           ["API 原生", "Anthropic memory / compaction",
            "笔记持久化、服务端压缩做成 API 原语；实测 agent 搜索 +39%"]],
          [0.16, 0.30, 0.54], body_size=9)
text_box(s1, M + Inches(0.2), TOP + Inches(2.62), CW - Inches(0.4), Inches(0.6), [
    [("◎ 前置契约：", 9, ACCENT2, True, False),
     ("在 system prompt 预声明标签语义（系统的话、与所附着消息无关、无需回复）。权威等级：", 9, MUTED, False, False),
     ("reminder > 用户 > 工具数据", 9, TEXT, True, False),
     ("，兼具抗注入价值。", 9, MUTED, False, False)],
])

# 右卡：四个结构性问题
RX = M + CW + Inches(0.25)
add_card(s1, RX, TOP, CW, CARD_H)
text_box(s1, RX + Inches(0.2), TOP + Inches(0.1), CW - Inches(0.4), Inches(0.3),
         [[("● 解决的四个结构性问题", 11.5, ACCENT2, True, False)]])
add_table(s1, RX + Inches(0.18), TOP + Inches(0.46), CW - Inches(0.36), Inches(2.5),
          ["问题", "表现", "reminder 解法"],
          [["注意力几何", "lost-in-the-middle，中段指令被淹没", "关键信息重注入 recency 高注意力区"],
           ["Context rot", "~80K token 后指令遵循度下降", "事件触发式规则刷新，在相关时刻重申"],
           ["状态漂移", "基于过期环境状态决策", "环境增量主动推送（文件变更、模式切换）"],
           ["缓存失效", "动态内容改写前缀，全量 miss", "一律 append 消息流尾部，前缀纹丝不动"]],
          [0.22, 0.39, 0.39], body_size=9)

# 底部横条卡：注入位置
BY = TOP + CARD_H + Inches(0.22)
BH = SLIDE_H - BY - Inches(0.42)
add_card(s1, M, BY, Inches(12.49), BH)
text_box(s1, M + Inches(0.2), BY + Inches(0.1), Inches(12.1), Inches(0.3),
         [[("● 注入位置：三种载体，一条铁律 —— 事件发生在哪个时机，就搭哪班车", 11.5, WARN, True, False)]])
add_table(s1, M + Inches(0.18), BY + Inches(0.46), Inches(7.1), BH - Inches(0.6),
          ["载体", "适用时机", "典型内容"],
          [["随 user message", "用户新一轮 query 到来", "时间/环境等动态上下文、模式切换、检索命中"],
           ["随 tool result", "loop 执行中发生事件（模型决策前唯一的新输入）", "文件外部被修改、安全提醒、隐私字段标识"],
           ["合成 user message", "无自然载体（会话恢复、压缩后、用户插话）", "上下文块、压缩后目标重申"]],
          [0.24, 0.40, 0.36], body_size=8.5)
text_box(s1, M + Inches(7.5), BY + Inches(0.44), Inches(4.8), BH - Inches(0.55), [
    [("不能注入的位置：", 9, TEXT, True, False),
     ("assistant message（伪造模型发言）；凭空的 tool result（必须对应真实 tool_use id，只能\u201c随\u201d不能\u201c造\u201d）；system prompt（缓存）。", 9, MUTED, False, False)],
    [("◎ 铁律：", 9, ACCENT2, True, False),
     ("只附着在 append 前沿（当前正要发出的新消息），绝不回头改历史消息 —— 否则其后 KV cache 全灭。", 9, TEXT, False, False)],
], space_after=5)
add_pageno(s1, "1 / 2")

# ══════════════════ 第 2 页：两个实践 ══════════════════
s2 = prs.slides.add_slide(blank)
set_bg(s2)
add_topbar(s2)

text_box(s2, M, Inches(0.18), Inches(3.6), Inches(0.75), [
    [("我们的两个实践", 21, TEXT, True, False)],
    [("覆盖 reminder 的两大载体", 10, MUTED, False, False)],
])
head2 = add_card(s2, Inches(4.15), Inches(0.16), Inches(8.77), Inches(0.78), fill=HEAD_BG)
text_box(s2, Inches(4.35), Inches(0.22), Inches(8.4), Inches(0.68), [
    [("tool result 侧约束", 9.5, MUTED, False, False),
     ("输出行为", 9.5, ACCENT2, True, False),
     ("（什么不能给用户看），user 侧供给", 9.5, MUTED, False, False),
     ("输入语境", 9.5, ACCENT2, True, False),
     ("（此刻用户在哪、现在几点）—— 前缀零改动，缓存全程命中。", 9.5, MUTED, False, False)],
], anchor=MSO_ANCHOR.MIDDLE, space_after=1)

TOP2 = Inches(1.08)
CARD_H2 = SLIDE_H - TOP2 - Inches(0.42)

# 左卡：实践一
add_card(s2, M, TOP2, CW, CARD_H2)
text_box(s2, M + Inches(0.2), TOP2 + Inches(0.12), CW - Inches(0.4), Inches(0.3),
         [[("实践一 · 随 tool result：隐私字段屏蔽", 11.5, ACCENT2, True, False)]])
text_box(s2, M + Inches(0.2), TOP2 + Inches(0.48), CW - Inches(0.4), Inches(0.85), [
    [("问题：", 9, TEXT, True, False),
     ("工具中间结果含隐私字段（标识符、号码等），推理需要用，但不应呈现给用户；逐个改造工具脱敏代价高且切断推理链。", 9, MUTED, False, False),
     ("做法：", 9, TEXT, True, False),
     ("框架层随结果附加 reminder，工具零改动：", 9, MUTED, False, False)],
])
add_code_block(s2, M + Inches(0.2), TOP2 + Inches(1.42), CW - Inches(0.4), Inches(1.95), [
    [("tool result:", RGBColor(0x8B, 0x9A, 0xB0), False)],
    [('{ "contacts": [ { "name": "张三",', CODE_TEXT, False)],
    [('    "phone": "13812341234", "uid": "u_8f3a92c1" } ] }', CODE_TEXT, False)],
    [("", CODE_TEXT, False)],
    [("<system-reminder>", ACCENT2, True)],
    [("以上结果中 phone / uid 为中间隐私字段，仅供你完成当前任务推理使用；", ACCENT2, False)],
    [("不得在给用户的回复中展示、复述或改写呈现。", ACCENT2, False)],
    [("</system-reminder>", ACCENT2, True)],
])
text_box(s2, M + Inches(0.2), TOP2 + Inches(3.55), CW - Inches(0.4), Inches(2.2), [
    [("• 数据可用、不外泄：", 9, TEXT, True, False),
     ("数据留在上下文（后续调用可能还要用 uid），只约束呈现行为", 9, MUTED, False, False)],
    [("• 就近生效：", 9, TEXT, True, False),
     ("约束紧贴数据出现的位置，遵循度远高于 system prompt 里泛化的\u201c注意隐私\u201d", 9, MUTED, False, False)],
    [("• 权威等级清晰：", 9, TEXT, True, False),
     ("即使用户要求\u201c把手机号发我\u201d，模型有依据拒绝或走脱敏路径", 9, MUTED, False, False)],
], space_after=5)

# 右卡：实践二
add_card(s2, RX, TOP2, CW, CARD_H2)
text_box(s2, RX + Inches(0.2), TOP2 + Inches(0.12), CW - Inches(0.4), Inches(0.3),
         [[("实践二 · 随 user message：时间与前台应用", 11.5, ACCENT, True, False)]])
text_box(s2, RX + Inches(0.2), TOP2 + Inches(0.48), CW - Inches(0.4), Inches(0.85), [
    [("问题：", 9, TEXT, True, False),
     ("用户指代依赖此刻设备状态——\u201c第4张照片\u201d是哪张、\u201c昨天\u201d是哪天；这些每轮都变，进 system prompt 等于每轮全量 cache miss。", 9, MUTED, False, False),
     ("做法：", 9, TEXT, True, False),
     ("结构化前缀随当轮 user message 注入：", 9, MUTED, False, False)],
])
add_code_block(s2, RX + Inches(0.2), TOP2 + Inches(1.42), CW - Inches(0.4), Inches(0.85), [
    [("user: ", RGBColor(0x8B, 0x9A, 0xB0), False),
     ("[用户输入时间：Tue 2026-08-25 15:43] [前台应用：图库]", ACCENT, True)],
    [("给我详细描述一下第4张照片", CODE_TEXT, False)],
])
text_box(s2, RX + Inches(0.2), TOP2 + Inches(2.45), CW - Inches(0.4), Inches(1.8), [
    [("• 指代消歧：", 9, TEXT, True, False),
     ("\u201c第4张照片\u201d → 前台是图库 → 当前视图第 4 张，无需反问、无需先调工具探测", 9, MUTED, False, False)],
    [("• 相对时间解析：", 9, TEXT, True, False),
     ("\u201c昨天 / 上周\u201d有了锚点；时间精确到分钟也不破坏缓存（在尾部，不在前缀）", 9, MUTED, False, False)],
    [("• 意图路由：", 9, TEXT, True, False),
     ("同一句话在不同前台应用下路由到不同 skill / 工具链", 9, MUTED, False, False)],
], space_after=5)
text_box(s2, RX + Inches(0.2), TOP2 + Inches(4.35), CW - Inches(0.4), Inches(1.5), [
    [("◎ 设计要点：", 9, ACCENT2, True, False),
     ("只随当轮、不回改历史（append-only）；[键：值] 固定模板、字段顺序确定；system prompt 预声明\u201c历史轮的值已过期，以最新一轮为准\u201d；只放高频消歧字段，大体积环境信息走主动加载。", 9, MUTED, False, False)],
])
add_pageno(s2, "2 / 2")

out = "/workspace/docs/reminder-discussion/reminder机制.pptx"
prs.save(out)
print("saved:", out)
